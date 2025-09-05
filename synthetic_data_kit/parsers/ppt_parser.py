"""
Parser for Microsoft PowerPoint presentations (PPTX).
"""
import os
import logging
from typing import Any, Optional

from synthetic_data_kit.parsers.base_parser import BaseParser

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger("ppt_parser")

class PPTParser(BaseParser):
    """Parser for Microsoft PowerPoint presentations (PPTX). Requires python-pptx."""
    
    def parse(self, file_path: str) -> str:
        """
        Parse a PPTX file and return its content.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            Extracted text content
        """
        logger.info(f"Parsing PPTX file: {file_path}")
        
        try:
            # Import pptx module
            from pptx import Presentation
            
            # Open the presentation
            prs = Presentation(file_path)
            
            # Extract text from slides
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Add slide number
                slide_text.append(f"Slide {i+1}")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # If there's more than just the slide number
                    slides_text.append("\n".join(slide_text))
            
            # Join slides with double newlines
            content = "\n\n".join(slides_text)
            
            # Clean the text
            content = self.clean_text(content)
            
            logger.info(f"Successfully parsed PPTX file ({len(content)} characters)")
            return content
        
        except ImportError:
            logger.error("PPTX parsing requires python-pptx. Please install it.")
            raise ImportError("PPTX parsing requires python-pptx to be installed.")
        
        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {str(e)}")
            raise